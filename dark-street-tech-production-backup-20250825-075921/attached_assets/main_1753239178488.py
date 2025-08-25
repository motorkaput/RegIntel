"""
Fetch Patterns AI - Document Analysis and Sentiment Analysis SaaS Platform
Copyright © 2025 Dark Street. All rights reserved.

This software is proprietary and confidential. Unauthorized copying, distribution,
or use of this software, via any medium, is strictly prohibited.
"""

from fastapi import FastAPI, File, UploadFile, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from typing import List
from pathlib import Path
import shutil
import os
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import docx
import pptx
import openpyxl
from transformers import pipeline
from keybert import KeyBERT
from openai import OpenAI
from dotenv import load_dotenv
import json
import re
from collections import Counter
from ai_service import AIService

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Initialize AI service and OpenAI client
ai_service = AIService()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY")) if os.getenv("OPENAI_API_KEY") else None

summary_store = []
classifier = pipeline("zero-shot-classification", model="facebook/bart-large-mnli")
sentiment_pipeline = pipeline("sentiment-analysis")
qa_pipeline = pipeline("question-answering")
kw_model = KeyBERT()

def extract_text(file_path):
    """Enhanced text extraction with better error handling"""
    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            return "\n".join([page.get_text() for page in doc])
        elif ext in [".jpg", ".jpeg", ".png"]:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        elif ext == ".docx":
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pptx":
            prs = pptx.Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(file_path)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            return "\n".join(text)
        elif ext in [".txt", ".text"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
    except Exception as e:
        return f"Error reading file {file_path.name}: {str(e)}"
    return ""

def get_contextual_sentiment_analysis(content: str, context_terms: list, num_documents: int = 1) -> dict:
    """Analyze sentiment for specific contexts within the document using OpenAI"""
    if not openai_client or not context_terms:
        return {}

    try:
        context_analysis = {}
        document_text = "document" if num_documents == 1 else "documents"

        for term in context_terms:
            prompt = f"""
            Analyze the sentiment specifically related to "{term}" in the following {document_text}.

            Document Content: {content[:2000]}

            Focus only on sentences, paragraphs, or sections that mention or relate to "{term}".

            Provide:
            1. Sentiment (positive/negative/neutral) specifically for this context
            2. Confidence score (0-1)
            3. Key phrases that influenced this sentiment
            4. Context summary (1-2 sentences about how {term} is discussed in the {document_text})
            5. Emotional tone (e.g., optimistic, concerned, frustrated, excited)

            Format as JSON: {{"sentiment": "positive/negative/neutral", "confidence": 0.85, "key_phrases": ["phrase1", "phrase2"], "context_summary": "summary", "emotional_tone": "tone"}}
            """

            response = openai_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
                temperature=0.2
            )

            try:
                content = response.choices[0].message.content.strip()
                # Try to extract JSON from the response if it's wrapped in markdown
                if content.startswith('```json'):
                    content = content.split('```json')[1].split('```')[0].strip()
                elif content.startswith('```'):
                    content = content.split('```')[1].split('```')[0].strip()

                analysis = json.loads(content)
                context_analysis[term] = analysis
            except (json.JSONDecodeError, IndexError) as e:
                print(f"JSON parsing failed for {term}: {e}")
                # Try to extract basic info from the raw response
                raw_content = response.choices[0].message.content

                # Simple sentiment extraction
                sentiment = "neutral"
                if "positive" in raw_content.lower():
                    sentiment = "positive"
                elif "negative" in raw_content.lower():
                    sentiment = "negative"

                context_analysis[term] = {
                    "sentiment": sentiment,
                    "confidence": 0.6,
                    "key_phrases": [],
                    "context_summary": raw_content[:200] + "..." if len(raw_content) > 200 else raw_content,
                    "emotional_tone": sentiment
                }

    except Exception as e:
        print(f"Contextual sentiment analysis failed: {e}")
        return {}

    return context_analysis

def transform_openai_analysis(openai_analysis, content, sentiment):
    """Transform OpenAI analysis to match frontend expectations"""
    if not openai_analysis:
        # Create basic structure from available data
        return {
            "sentiment": {
                "overall_sentiment": sentiment.get("label", "NEUTRAL").lower(),
                "confidence_score": sentiment.get("score", 0.5),
                "emotional_tone": ["neutral"],
                "business_impact": "Analysis not available",
                "recommendations": []
            },
            "business_insights": {
                "executive_summary": f"Document contains {len(content.split())} words with {sentiment.get('label', 'neutral').lower()} sentiment.",
                "key_insights": [],
                "opportunities": [],
                "risks": []
            },
            "compliance_assessment": {
                "risk_level": "medium",
                "regulatory_risks": [],
                "recommendations": []
            }
        }

    try:
        # Extract structured data from OpenAI analysis
        analysis = json.loads(openai_analysis) if isinstance(openai_analysis, str) else openai_analysis

        return {
            "sentiment": {
                "overall_sentiment": analysis.get("sentiment", {}).get("sentiment", sentiment.get("label", "NEUTRAL")).lower(),
                "confidence_score": analysis.get("sentiment", {}).get("confidence", sentiment.get("score", 0.5)),
                "emotional_tone": analysis.get("sentiment", {}).get("emotional_tone", ["neutral"]) if isinstance(analysis.get("sentiment", {}).get("emotional_tone"), list) else [analysis.get("sentiment", {}).get("emotional_tone", "neutral")],
                "business_impact": analysis.get("sentiment", {}).get("business_impact", "No specific business impact identified"),
                "recommendations": analysis.get("recommendations", [])[:3] if analysis.get("recommendations") else []
            },
            "business_insights": {
                "executive_summary": analysis.get("executive_summary", f"Document analysis: {len(content.split())} words with key themes identified."),
                "key_insights": analysis.get("key_insights", [])[:5] if analysis.get("key_insights") else [],
                "opportunities": analysis.get("opportunities", [])[:3] if analysis.get("opportunities") else [],
                "risks": analysis.get("risks", [])[:3] if analysis.get("risks") else []
            },
            "compliance_assessment": {
                "risk_level": analysis.get("risk_assessment", "medium") if analysis.get("risk_assessment") else "medium",
                "regulatory_risks": analysis.get("regulatory_risks", [])[:3] if analysis.get("regulatory_risks") else [],
                "recommendations": analysis.get("compliance_recommendations", [])[:3] if analysis.get("compliance_recommendations") else []
            }
        }
    except Exception as e:
        print(f"Error transforming OpenAI analysis: {e}")
        # Fallback structure
        return {
            "sentiment": {
                "overall_sentiment": sentiment.get("label", "NEUTRAL").lower(),
                "confidence_score": sentiment.get("score", 0.5),
                "emotional_tone": ["neutral"],
                "business_impact": "Analysis transformation failed",
                "recommendations": []
            },
            "business_insights": {
                "executive_summary": f"Document contains {len(content.split())} words. Analysis details unavailable.",
                "key_insights": [],
                "opportunities": [],
                "risks": []
            },
            "compliance_assessment": {
                "risk_level": "medium",
                "regulatory_risks": [],
                "recommendations": []
            }
        }

def get_openai_analysis(content, filename):
    """Enhanced analysis using OpenAI GPT-4"""
    if not openai_client:
        return None

    try:
        prompt = f"""
        Analyze the following document content and provide detailed insights:

        Document: {filename}
        Content: {content[:3000]}...

        Please provide:
        1. Document Summary: Create a comprehensive 3-4 sentence summary that captures the main purpose, key findings, and primary conclusions of the document. Focus on what the document is about, what it accomplishes, and its main value proposition.
        2. Key Themes (list 3-5 main themes)
        3. Sentiment Analysis (positive/negative/neutral with confidence)
        4. Document Type Classification
        5. Key Insights (3-5 actionable insights)
        6. Risk Assessment (if applicable)
        7. Recommendations (if applicable)

        For the summary, avoid generic phrases like "This document discusses" or "Document analysis shows". Instead, provide specific, meaningful content about what the document actually contains and accomplishes.

        Format as JSON with these keys: summary, themes, sentiment, document_type, insights, risks, recommendations
        """

        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",  # Using more cost-effective model
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1000,
            temperature=0.3
        )

        analysis_text = response.choices[0].message.content

        # Try to parse as JSON, fallback to structured text
        try:
            return json.loads(analysis_text)
        except:
            return {"raw_analysis": analysis_text}

    except Exception as e:
        print(f"OpenAI analysis failed: {e}")
        return None

def enhanced_keyword_extraction(content):
    """Enhanced keyword extraction with frequency analysis"""
    # Use KeyBERT for initial extraction
    keywords = kw_model.extract_keywords(content, keyphrase_ngram_range=(1, 3), stop_words="english", top_n=15)
    
    # Add frequency analysis
    words = re.findall(r'\b[a-zA-Z]{3,}\b', content.lower())
    word_freq = Counter(words)
    
    # Combine KeyBERT results with frequency analysis
    enhanced_keywords = []
    for kw, score in keywords:
        freq = word_freq.get(kw.lower(), 0)
        enhanced_keywords.append({
            "keyword": kw,
            "relevance_score": score,
            "frequency": freq,
            "combined_score": score * (1 + freq * 0.1)
        })
    
    return enhanced_keywords

@app.post("/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    summary_store.clear()
    results = []

    for file in files:
        file_path = Path(UPLOAD_DIR) / file.filename
        with open(file_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        content = extract_text(file_path)

        if not content or content.strip() == "":
            content = f"Empty or unreadable file: {file.filename}"

        # Enhanced keyword extraction
        enhanced_keywords = enhanced_keyword_extraction(content)

        # Traditional analysis (fallback)
        basic_keywords = [kw["keyword"] for kw in enhanced_keywords[:10]]
        sentiment = sentiment_pipeline(content[:512])[0] if content else {"label": "NEUTRAL", "score": 0.0}

        classification = classifier(
            content[:1024],
            candidate_labels=["business", "legal", "technical", "financial", "personal", "medical", "academic", "marketing"],
        )

        # AI Service enhanced analysis
        business_insights = await ai_service.generate_business_insights(content, file.filename)
        sentiment_analysis = await ai_service.analyze_sentiment(content)
        classification_result = await ai_service.classify_document(content, file.filename)

        # Transform AI service results to match frontend expectations
        ai_analysis = {
            "sentiment": {
                "overall_sentiment": sentiment_analysis.get("overall_sentiment", "neutral"),
                "confidence_score": sentiment_analysis.get("confidence_score", 0.5),
                "emotional_tone": sentiment_analysis.get("emotional_tone", ["neutral"]),
                "business_impact": sentiment_analysis.get("business_impact", "Analysis not available"),
                "recommendations": sentiment_analysis.get("recommendations", [])
            },
            "business_insights": {
                "executive_summary": business_insights.get("executive_summary", f"Document analysis: {len(content.split())} words with key themes identified."),
                "key_insights": business_insights.get("key_insights", []),
                "opportunities": business_insights.get("opportunities", []),
                "risks": business_insights.get("risks", [])
            },
            "compliance_assessment": {
                "risk_level": "medium",
                "regulatory_risks": [],
                "recommendations": []
            }
        }

        # Create result object for frontend compatibility
        result = {
            "filename": file.filename,
            "content": content,
            "keywords": basic_keywords,
            "enhanced_keywords": enhanced_keywords,
            "sentiment": sentiment,
            "classification": classification_result,
            "ai_analysis": ai_analysis,  # Frontend-compatible structure with AI service data
            "content_stats": {
                "word_count": len(content.split()),
                "char_count": len(content),
                "readability_score": calculate_readability(content)
            }
        }

        summary_store.append(result)
        results.append(result)

    return {
        "status": "success",
        "message": "Files uploaded and processed successfully",
        "results": results,
        "documents_processed": len(results)
    }

def calculate_readability(text):
    """Simple readability score calculation"""
    if not text:
        return 0
    
    sentences = len(re.findall(r'[.!?]+', text))
    words = len(text.split())
    
    if sentences == 0:
        return 0
    
    avg_sentence_length = words / sentences
    # Simplified readability score (lower is easier to read)
    return min(100, max(0, 100 - (avg_sentence_length * 2)))

@app.get("/summary")
async def get_summary():
    return JSONResponse(content=summary_store)

@app.post("/contextual-sentiment")
async def get_contextual_sentiment(request: Request):
    """Get contextual sentiment analysis for specific terms"""
    body = await request.json()
    context_terms = body.get("context_terms", [])

    if not context_terms:
        raise HTTPException(status_code=400, detail="No context terms provided")

    if not summary_store:
        raise HTTPException(status_code=400, detail="No documents processed")

    results = {}
    num_documents = len([doc for doc in summary_store if doc.get("content")])

    for doc in summary_store:
        if doc.get("content"):
            contextual_analysis = get_contextual_sentiment_analysis(doc["content"], context_terms, num_documents)
            if contextual_analysis:
                results[doc["filename"]] = contextual_analysis

    return {"contextual_sentiment": results}

def clean_answer_formatting(text: str) -> str:
    """
    Clean answer text to remove special characters and use simple bullets
    """
    if not text:
        return text

    # Remove markdown headers (### ## #)
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)

    # Remove bold/italic markdown (**text** *text*)
    text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)

    # Replace em-dashes and en-dashes with regular hyphens
    text = text.replace('—', '-').replace('–', '-')

    # Replace markdown bullets with simple bullets
    text = re.sub(r'^[\s]*[-*+]\s*', '• ', text, flags=re.MULTILINE)

    # Replace numbered lists with simple bullets
    text = re.sub(r'^[\s]*\d+\.\s*', '• ', text, flags=re.MULTILINE)

    # Clean up multiple spaces and normalize line breaks
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n\s*\n', '\n\n', text)

    return text.strip()

@app.post("/query")
async def ask_question(request: Request):
    body = await request.json()
    question = body.get("question", "")

    # Use AI service for enhanced Q&A
    try:
        # Prepare documents for AI service
        documents = [{"content": doc["content"], "filename": doc["filename"]}
                    for doc in summary_store if doc.get("content")]

        if not documents:
            return {
                "status": "error",
                "message": "No documents available for analysis"
            }

        # Use AI service smart Q&A
        qa_result = await ai_service.smart_qa(question, documents)

        # Clean the answer formatting
        clean_answer = clean_answer_formatting(qa_result.get("direct_answer", ""))

        return {
            "status": "success",
            "question": question,
            "answer": clean_answer,
            "source": "ai_service",
            "confidence": qa_result.get("confidence_level", 0.8),
            "supporting_evidence": qa_result.get("supporting_evidence", []),
            "source_documents": qa_result.get("source_documents", []),
            "related_insights": qa_result.get("related_insights", [])
        }

    except Exception as e:
        print(f"AI service Q&A failed: {e}")

        # Fallback to basic OpenAI if AI service fails
        if openai_client:
            try:
                combined_text = "\n".join([doc["content"] for doc in summary_store if doc["content"]])

                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You are an expert document analyst. Answer questions based on the provided documents. Provide clear, comprehensive answers using simple bullet points when needed. Avoid using special characters like ###, ***, em-dashes, or en-dashes."},
                        {"role": "user", "content": f"Documents: {combined_text[:4000]}\n\nQuestion: {question}"}
                    ],
                    max_tokens=1500,
                    temperature=0.2
                )

                raw_answer = response.choices[0].message.content
                clean_answer = clean_answer_formatting(raw_answer)

                return {
                    "status": "success",
                    "question": question,
                    "answer": clean_answer,
                    "source": "openai_fallback",
                    "confidence": 0.85
                }
            except Exception as e2:
                print(f"OpenAI fallback Q&A failed: {e2}")

        # Final fallback to traditional Q&A
        combined_text = "\n".join([doc["content"] for doc in summary_store if doc["content"]])
        answer = qa_pipeline(question=question, context=combined_text[:2000])
        clean_answer = clean_answer_formatting(answer["answer"])

        return {
            "status": "success",
            "question": question,
            "answer": clean_answer,
            "source": "transformers",
            "score": answer["score"]
        }

@app.get("/analytics")
async def get_analytics():
    """Enhanced analytics endpoint"""
    if not summary_store:
        return {"error": "No documents processed"}
    
    # Aggregate analytics
    total_docs = len(summary_store)
    total_words = sum(doc["content_stats"]["word_count"] for doc in summary_store)
    avg_readability = sum(doc["content_stats"]["readability_score"] for doc in summary_store) / total_docs
    
    # Sentiment distribution
    sentiment_dist = {}
    for doc in summary_store:
        sentiment = doc["sentiment"]["label"]
        sentiment_dist[sentiment] = sentiment_dist.get(sentiment, 0) + 1
    
    # Top keywords across all documents
    all_keywords = []
    for doc in summary_store:
        all_keywords.extend(doc.get("enhanced_keywords", []))
    
    keyword_scores = {}
    for kw_data in all_keywords:
        kw = kw_data["keyword"]
        score = kw_data["combined_score"]
        keyword_scores[kw] = keyword_scores.get(kw, 0) + score
    
    top_keywords = sorted(keyword_scores.items(), key=lambda x: x[1], reverse=True)[:20]
    
    return {
        "total_documents": total_docs,
        "total_words": total_words,
        "average_readability": round(avg_readability, 2),
        "sentiment_distribution": sentiment_dist,
        "top_keywords": [{"keyword": kw, "score": score} for kw, score in top_keywords],
        "document_types": [doc["classification"]["labels"][0] for doc in summary_store]
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
